from flask import Flask, request, jsonify
from flask_cors import CORS
from pptx import Presentation
import tempfile, os
from dotenv import load_dotenv
import google.generativeai as genai
load_dotenv()
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

app = Flask(__name__)
CORS(app)

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = []
    for idx, slide in enumerate(prs.slides, start=1):
        text_list = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                text_list.append(shape.text.strip())
        slide_text = "\n".join(text_list).strip()
        slides_text.append({"slide_index": idx, "text": slide_text})
    return slides_text

@app.route('/extract', methods=['POST'])
def extract():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400
    file = request.files['file']
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        file.save(tmp.name)
        slides_data = extract_text_from_pptx(tmp.name)
    os.remove(tmp.name)
    return jsonify(slides_data)

@app.route('/summarize', methods=['POST'])
def summarize():
    data = request.get_json()
    text = data.get('text', '')
    if not text:
        return jsonify({"error": "No text provided"}), 400

    # Set your Gemini API key here
    genai.configure(api_key=GEMINI_API_KEY)

    model = genai.GenerativeModel('gemini-2.0-flash')  # <-- updated model name
    response = model.generate_content(f"Summarize the following text:\n{text}")
    summary = response.text
    return jsonify({"summary": summary})

if __name__ == "__main__":
    app.run(debug=True)